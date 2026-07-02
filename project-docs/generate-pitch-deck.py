"""
OST Token — Satellite DePIN Pitch Deck Generator
Creates a professional PowerPoint presentation for SpaceX/Starlink partnership pitch.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0a, 0x0a, 0x14)
BG_CARD    = RGBColor(0x12, 0x12, 0x24)
ACCENT     = RGBColor(0x00, 0xd4, 0xff)
SUCCESS    = RGBColor(0x00, 0xff, 0x88)
WARNING    = RGBColor(0xff, 0xaa, 0x00)
TEXT_WHITE = RGBColor(0xff, 0xff, 0xff)
TEXT_DIM   = RGBColor(0x99, 0x99, 0xbb)
DANGER     = RGBColor(0xff, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16, color=TEXT_WHITE, spacing=1.2):
    """lines = list of (text, color, bold, size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text = item[0]
        c = item[1] if len(item) > 1 else color
        b = item[2] if len(item) > 2 else False
        sz = item[3] if len(item) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = "Segoe UI"
        p.space_after = Pt(font_size * (spacing - 1) * 2)
    return txBox


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent line at top
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
         "OST TOKEN", 60, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.8),
         "Satellite DePIN Micropayment Infrastructure", 32, TEXT_WHITE,
         bold=False, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
         "Universal digital cash that belongs to no country and serves every citizen.",
         20, TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.3), Inches(0.03), ACCENT)

add_text(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
         "Pitch Deck — Adventure", 22, TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.7), Inches(10), Inches(0.5),
         "Built on Solana  •  Token-2022  •  SpaceX Supplier Application",
         16, TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
         "nachogtavl-collab.github.io/ost-token", 14, ACCENT, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), DANGER)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "THE PROBLEM", 36, DANGER, bold=True)

# Problem cards
problems = [
    ("2.6 BILLION", "people have no internet access", "Starlink is connecting them — but with no native payment rail"),
    ("1.4 BILLION", "adults are unbanked", "They can receive internet but can't pay for it traditionally"),
    ("80+ COUNTRIES", "censor internet & payments", "Traditional payment processors comply with censorship orders"),
    ("4M+ STARLINK", "subscribers — zero P2P payments", "No way to trade bandwidth, share access, or settle micropayments"),
]

for i, (stat, desc, detail) in enumerate(problems):
    x = Inches(0.8 + (i % 2) * 6.1)
    y = Inches(1.5 + (i // 2) * 2.8)
    card = add_shape(slide, x, y, Inches(5.5), Inches(2.3), BG_CARD, ACCENT, Pt(1))
    add_text(slide, x + Inches(0.4), y + Inches(0.3), Inches(4.7), Inches(0.6),
             stat, 28, WARNING, bold=True)
    add_text(slide, x + Inches(0.4), y + Inches(0.9), Inches(4.7), Inches(0.4),
             desc, 18, TEXT_WHITE, bold=True)
    add_text(slide, x + Inches(0.4), y + Inches(1.5), Inches(4.7), Inches(0.6),
             detail, 14, TEXT_DIM)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), SUCCESS)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "THE SOLUTION — OST TOKEN", 36, SUCCESS, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "A Solana-based micropayment protocol purpose-built for satellite network operations",
         20, TEXT_DIM)

solutions = [
    ("🛰️  Bandwidth Micropayments", "Sub-cent, sub-second settlement for satellite bandwidth allocation. Starlink terminals tokenize and trade unused bandwidth."),
    ("🔒  Censorship-Resistant Payments", "Token-2022 confidential transfers with ElGamal + ZK proofs. Payments work where traditional rails are blocked."),
    ("📡  DePIN Node Rewards", "Automated on-chain rewards for ground stations, satellite operators, GPU compute, bandwidth, energy grid, and IoT nodes."),
    ("🌍  Disaster Response Network", "0.1% DAO fee funds satellite deployment for emergency connectivity. Self-sustaining humanitarian infrastructure."),
    ("🚀  Mars-Ready Protocol", "Offline tap-to-pay for environments with 3-22 minute communication delays. Multi-planetary by design."),
    ("👶  Grow Accounts", "Programmable savings for the next generation. On-chain accounts with release schedules. No bank required."),
]

for i, (title, desc) in enumerate(solutions):
    x = Inches(0.8 + (i % 3) * 4.1)
    y = Inches(2.0 + (i // 3) * 2.6)
    card = add_shape(slide, x, y, Inches(3.7), Inches(2.2), BG_CARD, ACCENT, Pt(1))
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.1), Inches(0.5),
             title, 16, ACCENT, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.1), Inches(1.2),
             desc, 12, TEXT_DIM)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — TECHNICAL SPECS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "TECHNICAL ARCHITECTURE", 36, ACCENT, bold=True)

# Left: specs table
specs = [
    ("Blockchain", "Solana (Token-2022)"),
    ("Transaction Cost", "< $0.0025"),
    ("Finality", "0.4 seconds"),
    ("Program Size", "572 KB BPF binary"),
    ("Instructions", "25 (Anchor framework)"),
    ("Token Supply", "1,000,000,000 OST"),
    ("Decimals", "9"),
    ("Privacy", "Confidential transfers (ElGamal + ZK)"),
    ("DAO Fee", "0.1% hardcoded, immutable"),
    ("SDK", "TypeScript, 22 instruction builders"),
    ("Pre-mine", "ZERO"),
    ("VC Funding", "ZERO"),
]

add_shape(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(5.8), BG_CARD, ACCENT, Pt(1))

for i, (key, val) in enumerate(specs):
    y = Inches(1.5 + i * 0.45)
    add_text(slide, Inches(0.8), y, Inches(2.8), Inches(0.4),
             key, 14, TEXT_DIM)
    add_text(slide, Inches(3.5), y, Inches(2.8), Inches(0.4),
             val, 14, TEXT_WHITE, bold=True)

# Right: DePIN categories
add_shape(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(5.8), BG_CARD, SUCCESS, Pt(1))
add_text(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.5),
         "DePIN REWARD CATEGORIES", 20, SUCCESS, bold=True)

categories = [
    ("🛰️  Satellite Nodes", "LEO orbit nodes validating transactions"),
    ("📡  Ground Stations", "Uplink/downlink infrastructure operators"),
    ("💻  GPU Compute", "Distributed processing power for the network"),
    ("🌐  Bandwidth Sharing", "P2P internet sharing via Starlink terminals"),
    ("🔋  Energy Grid", "Renewable energy for ground stations"),
    ("⚡  IoT Devices", "Sensor and edge device contributions"),
]

for i, (cat, desc) in enumerate(categories):
    y = Inches(2.3 + i * 0.8)
    add_text(slide, Inches(7.3), y, Inches(5.2), Inches(0.35),
             cat, 16, ACCENT, bold=True)
    add_text(slide, Inches(7.3), y + Inches(0.35), Inches(5.2), Inches(0.35),
             desc, 12, TEXT_DIM)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — SPACEX ALIGNMENT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "WHY SPACEX × OST", 36, ACCENT, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Every SpaceX mission objective has a payment layer requirement that OST fulfills",
         18, TEXT_DIM)

alignments = [
    ("STARLINK BILLING", "Currently: flat monthly rate\nWith OST: Usage-based micro-billing, P2P bandwidth trading, fair-share pricing for shared terminals",
     "4M+ subscribers, growing 50%/year", ACCENT),
    ("CENSORED REGIONS", "Problem: Payment processors block transactions in 80+ countries\nWith OST: Confidential transfers bypass censorship, users pay for Starlink access directly",
     "Largest untapped market for Starlink", WARNING),
    ("DISASTER RESPONSE", "Currently: FEMA contracts, government-funded\nWith OST: Self-funding via 0.1% DAO treasury, autonomous deployment, no bureaucratic delay",
     "Hurricane/earthquake rapid connectivity", SUCCESS),
    ("MARS COLONIZATION", "Problem: 3-22 min light delay makes real-time payments impossible\nWith OST: Offline tap-to-pay, local finality, eventual Earth settlement",
     "Required for any Mars economy", DANGER),
]

for i, (title, desc, stat, color) in enumerate(alignments):
    x = Inches(0.8 + (i % 2) * 6.1)
    y = Inches(2.0 + (i // 2) * 2.6)
    card = add_shape(slide, x, y, Inches(5.5), Inches(2.2), BG_CARD, color, Pt(1.5))
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(4.9), Inches(0.4),
             title, 18, color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(4.9), Inches(0.9),
             desc, 12, TEXT_DIM)
    add_text(slide, x + Inches(0.3), y + Inches(1.7), Inches(4.9), Inches(0.3),
             stat, 13, TEXT_WHITE, bold=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — PROVEN PRECEDENTS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "PROVEN PATH — PRECEDENTS", 36, ACCENT, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Others have done parts of this. We combine all of them.", 18, TEXT_DIM)

precedents = [
    ("SPACECOIN", "Launched a 3U CubeSat via SpaceX Rideshare for ~$275K. Proved that crypto projects can literally put satellites in orbit.",
     "What they proved: Crypto can launch hardware into space\nWhat OST adds: Payment protocol + DePIN rewards, not just a satellite", ACCENT),
    ("WISeSat / WISeKey", "IoT security satellites launched on SpaceX. Enterprise-grade secure communications from LEO orbit.",
     "What they proved: Enterprise satellite IoT is viable via SpaceX\nWhat OST adds: Decentralized payment layer for IoT settlements", WARNING),
    ("X CORP (TWITTER)", "Building payments infrastructure. Elon Musk's vision: X as 'everything app' with integrated payments.",
     "What they proved: Musk wants payment rails for his platforms\nWhat OST adds: Censorship-resistant rail for X in blocked regions", SUCCESS),
]

for i, (name, desc, detail, color) in enumerate(precedents):
    y = Inches(2.0 + i * 1.7)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.5), BG_CARD, color, Pt(1))
    add_text(slide, Inches(1.2), y + Inches(0.15), Inches(2), Inches(0.4),
             name, 20, color, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.55), Inches(4.5), Inches(0.8),
             desc, 12, TEXT_DIM)
    add_text(slide, Inches(6.2), y + Inches(0.2), Inches(5.8), Inches(1.1),
             detail, 12, TEXT_WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — TREASURY & TOKENOMICS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), SUCCESS)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "TOKENOMICS & TREASURY", 36, SUCCESS, bold=True)

# Token stats
token_stats = [
    ("Total Supply", "1,000,000,000 OST"),
    ("Pre-mine", "ZERO"),
    ("VC Allocation", "ZERO"),
    ("Team Allocation", "ZERO"),
    ("DAO Fee", "0.1% per transaction"),
    ("Fee Destination", "DAO Treasury (PDA)"),
    ("Fee Mutability", "IMMUTABLE (hardcoded)"),
    ("Treasury Use", "Fund satellite infrastructure"),
]

add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(4.2), BG_CARD, SUCCESS, Pt(1))
add_text(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
         "TOKEN PARAMETERS", 18, SUCCESS, bold=True)

for i, (key, val) in enumerate(token_stats):
    y = Inches(2.1 + i * 0.45)
    add_text(slide, Inches(0.8), y, Inches(2.5), Inches(0.4),
             key, 14, TEXT_DIM)
    add_text(slide, Inches(3.2), y, Inches(2.5), Inches(0.4),
             val, 14, TEXT_WHITE, bold=True)

# Flywheel
add_shape(slide, Inches(6.5), Inches(1.3), Inches(6.3), Inches(4.2), BG_CARD, ACCENT, Pt(1))
add_text(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.4),
         "SELF-FUNDING FLYWHEEL", 18, ACCENT, bold=True)

flywheel = [
    "1.  Users transact with OST",
    "2.  0.1% fee automatically goes to DAO Treasury",
    "3.  Treasury funds satellite hardware & deployment",
    "4.  More satellites = more connectivity",
    "5.  More connectivity = more users",
    "6.  More users = more transactions",
    "7.  REPEAT → self-sustaining cycle",
]

for i, step in enumerate(flywheel):
    y = Inches(2.2 + i * 0.45)
    color = ACCENT if i < 6 else SUCCESS
    add_text(slide, Inches(6.8), y, Inches(5.7), Inches(0.4),
             step, 14, color, bold=(i == 6))

# Bottom note
add_shape(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), BG_CARD, WARNING, Pt(1))
add_text(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1),
         "Fair Launch Principles: Zero pre-mine means no insider advantage. Zero VC means no external pressure. "
         "The 0.1% DAO fee is the ONLY revenue mechanism, and it's hardcoded — no governance vote can change it. "
         "This ensures permanent alignment: the protocol only profits when users transact.",
         14, TEXT_DIM)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — ROADMAP
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "ROADMAP", 36, ACCENT, bold=True)

phases = [
    ("PHASE 1\nENTRY POINT", "Week 1  •  IN PROGRESS", [
        ("✅", "Live website & GitHub (open-source)"),
        ("✅", "DePIN faucet & treasury (0.1% fee on devnet)"),
        ("✅", "SpaceX supplier application submitted"),
        ("✅", "X visibility campaign launched"),
        ("✅", "Pitch deck created"),
    ], SUCCESS),
    ("PHASE 2\nPROVE IT", "Month 0–3", [
        ("▶️", "Mainnet launch & first 1,000 users"),
        ("▶️", "DePIN rewards flowing (real OST)"),
        ("▶️", "Daily X content, 10K+ impressions/week"),
    ], ACCENT),
    ("PHASE 3\nPARTNERSHIP", "Month 3–6", [
        ("▶️", "SpaceX supplier review (traction data)"),
        ("▶️", "Starlink Enterprise program pilot"),
        ("▶️", "X Corp payments integration"),
    ], WARNING),
    ("PHASE 4\nSPACE LAUNCH", "Month 6–18", [
        ("▶️", "SpaceX Rideshare CubeSat (~$275K)"),
        ("▶️", "OST mesh network — uncensored internet"),
        ("▶️", "Multi-planetary payment infrastructure"),
    ], DANGER),
]

for i, (title, timeline, items, color) in enumerate(phases):
    x = Inches(0.4 + i * 3.2)
    phase_card = add_shape(slide, x, Inches(1.3), Inches(2.9), Inches(5.8), BG_CARD, color, Pt(1.5))

    add_text(slide, x + Inches(0.15), Inches(1.5), Inches(2.6), Inches(0.7),
             title, 14, color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), Inches(2.2), Inches(2.6), Inches(0.3),
             timeline, 10, TEXT_DIM, alignment=PP_ALIGN.CENTER)

    for j, (icon, text) in enumerate(items):
        y = Inches(2.7 + j * 0.55)
        item_color = SUCCESS if icon == "✅" else TEXT_DIM
        add_text(slide, x + Inches(0.15), y, Inches(2.6), Inches(0.5),
                 f"{icon}  {text}", 11, item_color)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — PROGRESS / WHAT'S BUILT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), SUCCESS)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "WHAT'S ALREADY BUILT — 69% COMPLETE", 36, SUCCESS, bold=True)

# Progress bar
add_shape(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.35), BG_CARD)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(11.7 * 0.69), Inches(0.35), SUCCESS)
add_text(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.3),
         "11 of 16 milestones complete", 12, TEXT_DIM, alignment=PP_ALIGN.CENTER)

done_items = [
    "Anchor Smart Contract (25 instructions)",
    "TypeScript SDK (22 builders)",
    "Website & Frontend (3D globe, i18n, wallet)",
    "Security Audit (internal review)",
    "Program Compiled (572KB BPF)",
    "Token Mint (Token-2022, 1B supply)",
    "Devnet Deployment (program + mint + treasury)",
    "wOST Wrapper (SPL for AMM/DEX)",
    "Token Metadata On-Chain",
    "Raydium Liquidity Pool (wOST/SOL)",
    "DAO Treasury Wallet (PDA)",
]

pending_items = [
    "Mainnet Deployment (~2.5 SOL rent)",
    "Mainnet Liquidity Seeding",
    "Jupiter Strict List",
    "CoinGecko / CoinMarketCap",
    "External Security Audit",
]

add_shape(slide, Inches(0.5), Inches(2.2), Inches(7.5), Inches(4.8), BG_CARD, SUCCESS, Pt(1))
add_text(slide, Inches(0.8), Inches(2.4), Inches(7), Inches(0.4),
         "✅  COMPLETED (11)", 16, SUCCESS, bold=True)

for i, item in enumerate(done_items):
    col = i // 6
    row = i % 6
    x = Inches(0.8 + col * 3.5)
    y = Inches(3.0 + row * 0.6)
    add_text(slide, x, y, Inches(3.3), Inches(0.5),
             f"✅  {item}", 11, TEXT_WHITE)

add_shape(slide, Inches(8.3), Inches(2.2), Inches(4.5), Inches(4.8), BG_CARD, WARNING, Pt(1))
add_text(slide, Inches(8.6), Inches(2.4), Inches(4), Inches(0.4),
         "🟡  PENDING (5)", 16, WARNING, bold=True)

for i, item in enumerate(pending_items):
    y = Inches(3.0 + i * 0.7)
    add_text(slide, Inches(8.6), y, Inches(4), Inches(0.6),
             f"▶️  {item}", 12, TEXT_DIM)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — THE ASK
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "THE ASK", 36, ACCENT, bold=True)

add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
         "We're not asking for funding. We're asking for partnership.", 24, TEXT_WHITE, bold=True)

asks = [
    ("1. SUPPLIER STATUS", "Accept Adventure as a SpaceX supplier for satellite payment infrastructure.",
     "This validates OST as a recognized technology partner in the SpaceX ecosystem.", ACCENT),
    ("2. STARLINK API ACCESS", "Provide sandbox access to Starlink terminal APIs for bandwidth tokenization.",
     "We'll build the micropayment integration at zero cost to SpaceX.", SUCCESS),
    ("3. RIDESHARE CONSIDERATION", "When traction warrants it (~1K+ users, treasury > $50K), consider OST for a CubeSat rideshare slot.",
     "~$275K for 50kg CubeSat on Transporter mission. Treasury-funded.", WARNING),
]

for i, (title, desc, detail, color) in enumerate(asks):
    y = Inches(2.3 + i * 1.6)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.3), BG_CARD, color, Pt(1.5))
    add_text(slide, Inches(1.2), y + Inches(0.15), Inches(3), Inches(0.4),
             title, 20, color, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.55), Inches(5), Inches(0.6),
             desc, 14, TEXT_WHITE)
    add_text(slide, Inches(6.5), y + Inches(0.3), Inches(5.5), Inches(0.7),
             detail, 13, TEXT_DIM)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — CONTACT / CLOSING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(1.5), Inches(1), Inches(10), Inches(0.8),
         "OST TOKEN", 48, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2), Inches(10), Inches(0.6),
         '"Universal digital cash that belongs to no country\nand serves every citizen."',
         22, TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.0), Inches(2.3), Inches(0.03), ACCENT)

# Contact card
card = add_shape(slide, Inches(3.5), Inches(3.5), Inches(6.3), Inches(2.8), BG_CARD, ACCENT, Pt(1))

contact_lines = [
    ("ADVENTURE", ACCENT, True, 22),
    ("", TEXT_DIM, False, 8),
    ("Website:  nachogtavl-collab.github.io/ost-token", TEXT_WHITE, False, 14),
    ("GitHub:   github.com/nachogtavl-collab/ost-token", TEXT_WHITE, False, 14),
    ("Program:  J2jiS296YWVie1Sopb4SxcM3aJnP9aAwe6aLDhCqvGXY", TEXT_DIM, False, 11),
    ("OST Mint: 383pTzoZ8Gp83dzk23ZnvLcfX2Sq32TAGN48CMQu2pAJ", TEXT_DIM, False, 11),
    ("Network:  Solana Devnet (mainnet-ready)", TEXT_DIM, False, 11),
]

add_multiline(slide, Inches(3.8), Inches(3.7), Inches(5.7), Inches(2.5), contact_lines)

add_text(slide, Inches(1.5), Inches(6.7), Inches(10), Inches(0.4),
         "Built by one developer  •  From Mexico  •  Incorporated in Panama  •  Zero VC  •  Zero pre-mine",
         13, TEXT_DIM, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────
output_path = os.path.join(
    r"C:\Users\neyma\OneDrive\Desktop\New folder\ost-token\docs",
    "OST-Pitch-Deck.pptx"
)
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
